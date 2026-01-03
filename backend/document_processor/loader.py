import os
from pathlib import Path
from typing import List, Dict, Any
from langchain_core.documents import Document
from docx import Document as DocxDocument
from pptx import Presentation
import PyPDF2


class DocumentLoader:
    """Loads documents from Word, PowerPoint, PDF, and text files"""
    
    SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".pdf", ".txt"}
    
    @staticmethod
    def load_word_document(file_path: str) -> List[Document]:
        """Load content from Word documents (.docx)"""
        documents = []
        try:
            doc = DocxDocument(file_path)
            full_text = "\n".join([para.text for para in doc.paragraphs])
            
            # Extract tables if present
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        full_text += "\n" + cell.text
            
            file_name = Path(file_path).name
            documents.append(
                Document(
                    page_content=full_text,
                    metadata={
                        "source": file_name,
                        "file_type": "docx",
                        "file_path": file_path
                    }
                )
            )
        except Exception as e:
            print(f"Error loading Word document {file_path}: {e}")
        
        return documents
    
    @staticmethod
    def load_powerpoint_document(file_path: str) -> List[Document]:
        """Load content from PowerPoint documents (.pptx)"""
        documents = []
        try:
            prs = Presentation(file_path)
            file_name = Path(file_path).name
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text = f"Slide {slide_idx + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    documents.append(
                        Document(
                            page_content=slide_text,
                            metadata={
                                "source": file_name,
                                "file_type": "pptx",
                                "file_path": file_path,
                                "slide_number": slide_idx + 1
                            }
                        )
                    )
        except Exception as e:
            print(f"Error loading PowerPoint document {file_path}: {e}")
        
        return documents
    
    @staticmethod
    def load_pdf_document(file_path: str) -> List[Document]:
        """Load content from PDF documents"""
        documents = []
        try:
            file_name = Path(file_path).name
            with open(file_path, "rb") as pdf_file:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                
                for page_idx, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    
                    documents.append(
                        Document(
                            page_content=page_text,
                            metadata={
                                "source": file_name,
                                "file_type": "pdf",
                                "file_path": file_path,
                                "page_number": page_idx + 1
                            }
                        )
                    )
        except Exception as e:
            print(f"Error loading PDF document {file_path}: {e}")
        
        return documents
    
    @staticmethod
    def load_text_document(file_path: str) -> List[Document]:
        """Load content from text files"""
        documents = []
        try:
            file_name = Path(file_path).name
            with open(file_path, "r", encoding="utf-8") as txt_file:
                content = txt_file.read()
                
            documents.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": file_name,
                        "file_type": "txt",
                        "file_path": file_path
                    }
                )
            )
        except Exception as e:
            print(f"Error loading text document {file_path}: {e}")
        
        return documents
    
    @classmethod
    def load_document(cls, file_path: str) -> List[Document]:
        """Load a document based on its file extension"""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext == ".docx":
            return cls.load_word_document(file_path)
        elif file_ext == ".pptx":
            return cls.load_powerpoint_document(file_path)
        elif file_ext == ".pdf":
            return cls.load_pdf_document(file_path)
        elif file_ext == ".txt":
            return cls.load_text_document(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}. Supported types: {cls.SUPPORTED_EXTENSIONS}")
    
    @classmethod
    def load_documents_from_directory(cls, directory_path: str) -> List[Document]:
        """Load all supported documents from a directory"""
        documents = []
        directory = Path(directory_path)
        
        for file_path in directory.iterdir():
            if file_path.suffix.lower() in cls.SUPPORTED_EXTENSIONS:
                try:
                    loaded_docs = cls.load_document(str(file_path))
                    documents.extend(loaded_docs)
                except Exception as e:
                    print(f"Failed to load {file_path}: {e}")
        
        return documents
